"""Color and annotate the user-formatted h24 cascade PowerPoint in place.

The source presentation's box geometry and native connectors are preserved.
Outputs are written to a new PPTX so the supplied presentation remains intact.
"""

from __future__ import annotations

import argparse
from copy import deepcopy
from pathlib import Path
from typing import Mapping

import matplotlib
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt

from empirical_outcomes import load_empirical_outcomes
from experiments.m3n_vc._legacy_h24_report_plot import (
    DEFAULT_APPROXIMATE_REPORT,
    DEFAULT_BASELINE_REPORT,
    DEFAULT_BRUTE_FORCE_REPORT,
    DEFAULT_BRUTE_FORCE_RESULTS,
    DEFAULT_OUTCOMES,
    _holdout_evaluator,
    load_methods,
    occurrence_route_counts,
)
from threshold_optimizer import split_empirical_outcomes


DEFAULT_SOURCE = Path(r"C:\Users\TheSandwichCoder\Downloads\Cascade 096 layouts.pptx")
DEFAULT_OUTPUT = Path(
    r"C:\Users\TheSandwichCoder\Downloads\Cascade 096 layouts refined.pptx"
)

# PowerPoint shape indexes (zero-based) mapped to saved cascade locations.
BASELINE_LOCATIONS = {
    0: "initial[0]",
    1: "initial[1]",
    2: "initial[2]",
    3: "initial[3]",
    4: "initial[4]",
    5: "specialized[K0:suv][0]",
    7: "specialized[K0:suv][1]",
    16: "specialized[K0:suv][2]",
    17: "specialized[K0:suv][3]",
    6: "specialized[K0:coupe][0]",
    8: "specialized[K0:coupe][1]",
    18: "specialized[K0:coupe][2]",
    19: "specialized[K0:coupe][3]",
    11: "specialized[K1:suv][0]",
    9: "specialized[K1:coupe][0]",
    10: "specialized[K1:coupe][1]",
}
LINEAR_LOCATIONS = {0: "initial[0]", 1: "initial[1]"}
BRUTE_LOCATIONS = {
    0: "initial[0]",
    1: "initial[1]",
    7: "initial[2]",
    9: "initial[3]",
    2: "specialized[K0:suv][0]",
    11: "specialized[K0:suv][1]",
    3: "specialized[K0:coupe][0]",
    13: "specialized[K0:coupe][1]",
}
APPROXIMATE_LOCATIONS = {
    7: "initial[0]",
    0: "initial[1]",
    1: "initial[2]",
    9: "initial[3]",
    2: "specialized[K0:suv][0]",
    3: "specialized[K0:coupe][0]",
    11: "specialized[K0:coupe][1]",  # added below
}


def _rgb_from_share(share: float) -> RGBColor:
    red, green, blue, _ = matplotlib.colormaps["Blues"](share)
    return RGBColor(round(255 * red), round(255 * green), round(255 * blue))


def _set_run_style(run, *, size: float, color: RGBColor, bold: bool = False) -> None:
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _set_box_text(
    shape,
    candidate_id: str,
    threshold: float | None,
    share: float,
) -> None:
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.margin_left = Pt(2)
    text_frame.margin_right = Pt(2)
    text_frame.margin_top = Pt(1)
    text_frame.margin_bottom = Pt(1)
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.word_wrap = False
    color = RGBColor(255, 255, 255) if share >= 0.58 else RGBColor(20, 20, 20)

    lines = ["Kdet" if candidate_id == "detector" else candidate_id]
    if threshold is not None:
        lines.append(f"({threshold:.3f})")
    lines.append(f"{share:.1%} routes")
    for index, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        paragraph.line_spacing = 0.9
        run = paragraph.add_run()
        run.text = line
        _set_run_style(
            run,
            size=12.0 if index == 0 else 9.0,
            color=color,
            bold=index == 0,
        )


def _add_branch_label(slide, text: str, left: int, top: int) -> None:
    box = slide.shapes.add_textbox(left, top, 500000, 180000)
    box.fill.background()
    box.line.fill.background()
    frame = box.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = Pt(0)
    frame.margin_top = frame.margin_bottom = Pt(0)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    _set_run_style(run, size=8.0, color=RGBColor(75, 75, 75), bold=False)


def _candidate_at_location(evaluator, location: str) -> tuple[str, str]:
    if location.startswith("initial["):
        index = int(location.removeprefix("initial[").removesuffix("]"))
        candidate_id = evaluator.cascade.initial[index]
    else:
        body = location.removeprefix("specialized[")
        route, index_text = body.rsplit("][", maxsplit=1)
        index = int(index_text.removesuffix("]"))
        router_id, group = route.split(":", maxsplit=1)
        candidate_id = evaluator.cascade.specialized[(router_id, group)][index]
    occurrence = (
        f"detector@{location}"
        if candidate_id == evaluator.detector_id
        else evaluator._slot_by_location[location]  # noqa: SLF001
    )
    return candidate_id, occurrence


def _annotate_slide(
    slide,
    method,
    holdout_payload: Mapping[str, object],
    locations: Mapping[int, str],
) -> None:
    evaluator = _holdout_evaluator(method.layout, holdout_payload)
    thresholds = {
        str(key): float(value) for key, value in method.holdout["thresholds"].items()
    }
    normalized = evaluator._normalise_thresholds(thresholds)  # noqa: SLF001
    counts = occurrence_route_counts(evaluator, thresholds)

    for shape_index, location in locations.items():
        shape = slide.shapes[shape_index]
        candidate_id, occurrence = _candidate_at_location(evaluator, location)
        share = counts.get(occurrence, 0) / evaluator.sample_count
        threshold = None if candidate_id == evaluator.detector_id else normalized[occurrence]
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb_from_share(share)
        _set_box_text(shape, candidate_id, threshold, share)


def _add_missing_approximate_detector(slide) -> None:
    """Add the Kdet fallback omitted after the lower K6 branch."""

    k6 = slide.shapes[3]
    main_detector = slide.shapes[9]
    detector_element = deepcopy(main_detector.element)
    slide.shapes._spTree.insert_element_before(  # noqa: SLF001
        detector_element, "p:extLst"
    )
    detector = slide.shapes[-1]
    detector.left = main_detector.left
    detector.top = k6.top
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        k6.left + k6.width,
        k6.top + k6.height // 2,
        detector.left,
        detector.top + detector.height // 2,
    )
    connector.line.color.rgb = RGBColor(89, 89, 89)
    connector.line.width = Pt(0.75)


def refine_powerpoint(source: Path = DEFAULT_SOURCE, output: Path = DEFAULT_OUTPUT) -> None:
    if not source.is_file():
        raise FileNotFoundError(source)
    presentation = Presentation(source)
    if len(presentation.slides) != 5:
        raise ValueError("Expected exactly five cascade-layout slides.")

    methods = {method.key: method for method in load_methods()}
    payload = load_empirical_outcomes(DEFAULT_OUTCOMES)
    _, holdout_payload, _ = split_empirical_outcomes(
        payload,
        holdout_fraction=0.20,
        split_strategy="blocked_per_run",
        random_seed=0,
    )

    _annotate_slide(
        presentation.slides[0], methods["baseline"], holdout_payload, BASELINE_LOCATIONS
    )
    _annotate_slide(
        presentation.slides[1],
        methods["baseline_annealed"],
        holdout_payload,
        BASELINE_LOCATIONS,
    )
    _annotate_slide(
        presentation.slides[2], methods["linear_k3"], holdout_payload, LINEAR_LOCATIONS
    )
    _annotate_slide(
        presentation.slides[3], methods["brute_force"], holdout_payload, BRUTE_LOCATIONS
    )

    # Match the saved approximate topology while retaining the user's positions.
    presentation.slides[4].shapes[1].text = "K2"
    _add_missing_approximate_detector(presentation.slides[4])
    _annotate_slide(
        presentation.slides[4],
        methods["approximate_joint"],
        holdout_payload,
        APPROXIMATE_LOCATIONS,
    )

    # Small branch labels, placed in the existing whitespace beside dashed edges.
    for slide_index in (0, 1):
        slide = presentation.slides[slide_index]
        _add_branch_label(slide, "SUV", 1570000, 1420000)
        _add_branch_label(slide, "COUPE", 1570000, 3440000)
        _add_branch_label(slide, "SUV", 5660000, 1840000)
        _add_branch_label(slide, "COUPE", 5660000, 3020000)
    _add_branch_label(presentation.slides[3], "SUV", 2400000, 1650000)
    _add_branch_label(presentation.slides[3], "COUPE", 2400000, 3280000)
    _add_branch_label(presentation.slides[4], "SUV", 4410000, 1660000)
    _add_branch_label(presentation.slides[4], "COUPE", 4410000, 3300000)

    output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output)
    print(f"Wrote {output}")


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--source", type=Path, default=DEFAULT_SOURCE)
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT)
    args = parser.parse_args()
    refine_powerpoint(args.source, args.output)


if __name__ == "__main__":
    main()
